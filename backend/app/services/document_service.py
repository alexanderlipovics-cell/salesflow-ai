"""
Document Generation Service - PDF, PPTX, XLSX, DOCX
"""

import io
from datetime import datetime
from typing import List, Dict, Any
import logging

logger = logging.getLogger(__name__)


def generate_pdf(title: str, content: List[Dict[str, Any]]) -> bytes:
    """
    Generiert ein PDF.
    
    content = [
        {"type": "heading", "text": "Kapitel 1", "level": 1},
        {"type": "paragraph", "text": "Lorem ipsum..."},
        {"type": "bullet_list", "items": ["Item 1", "Item 2"]},
        {"type": "table", "headers": ["A", "B"], "rows": [["1", "2"]]},
    ]
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    
    styles.add(ParagraphStyle(
        name='CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.HexColor('#0891b2')
    ))
    
    story = []
    story.append(Paragraph(title, styles['CustomTitle']))
    story.append(Spacer(1, 20))
    
    for item in content:
        item_type = item.get("type")
        
        if item_type == "heading":
            level = item.get("level", 1)
            style = f"Heading{min(level, 3)}"
            story.append(Paragraph(item["text"], styles[style]))
            story.append(Spacer(1, 10))
            
        elif item_type == "paragraph":
            story.append(Paragraph(item["text"], styles["Normal"]))
            story.append(Spacer(1, 10))
            
        elif item_type == "bullet_list":
            items = [ListItem(Paragraph(i, styles["Normal"])) for i in item["items"]]
            story.append(ListFlowable(items, bulletType='bullet'))
            story.append(Spacer(1, 10))
            
        elif item_type == "table":
            headers = item.get("headers", [])
            rows = item.get("rows", [])
            data = [headers] + rows
            
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0891b2')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f0f9ff')),
                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e0e0e0')),
            ]))
            story.append(table)
            story.append(Spacer(1, 15))
    
    story.append(Spacer(1, 30))
    story.append(Paragraph(
        f"Erstellt mit AlSales AI • {datetime.now().strftime('%d.%m.%Y')}",
        ParagraphStyle(name='Footer', fontSize=8, textColor=colors.gray)
    ))
    
    doc.build(story)
    buffer.seek(0)
    return buffer.read()


def generate_pptx(title: str, slides: List[Dict[str, Any]]) -> bytes:
    """
    Generiert eine PowerPoint.
    
    slides = [
        {"title": "Slide 1", "content": ["Bullet 1", "Bullet 2"]},
        {"title": "Slide 2", "content": "Paragraph text"},
    ]
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(8, 145, 178)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Erstellt mit AlSales AI • {datetime.now().strftime('%d.%m.%Y')}"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # Content Slides
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        
        content = slide_data.get("content", [])
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        if isinstance(content, list):
            for i, item in enumerate(content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(20)
        else:
            p = tf.paragraphs[0]
            p.text = str(content)
            p.font.size = Pt(18)
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def generate_xlsx(title: str, sheets: List[Dict[str, Any]]) -> bytes:
    """
    Generiert eine Excel-Datei.
    
    sheets = [
        {
            "name": "Sheet1",
            "headers": ["Name", "Alter", "Stadt"],
            "rows": [["Max", 28, "Berlin"], ["Anna", 34, "München"]]
        }
    ]
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    
    wb = Workbook()
    wb.remove(wb.active)
    
    header_fill = PatternFill(start_color="0891b2", end_color="0891b2", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )
    
    for sheet_data in sheets:
        ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
        
        headers = sheet_data.get("headers", [])
        rows = sheet_data.get("rows", [])
        
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border
        
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal='center')
        
        for col in ws.columns:
            max_length = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = max_length + 2
    
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.read()


def generate_docx(title: str, content: List[Dict[str, Any]]) -> bytes:
    """
    Generiert ein Word-Dokument.
    """
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = Document()
    
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for item in content:
        item_type = item.get("type")
        
        if item_type == "heading":
            doc.add_heading(item["text"], item.get("level", 1))
        elif item_type == "paragraph":
            doc.add_paragraph(item["text"])
        elif item_type == "bullet_list":
            for bullet in item["items"]:
                doc.add_paragraph(bullet, style='List Bullet')
        elif item_type == "table":
            headers = item.get("headers", [])
            rows = item.get("rows", [])
            table = doc.add_table(rows=1 + len(rows), cols=len(headers))
            table.style = 'Table Grid'
            for i, header in enumerate(headers):
                table.rows[0].cells[i].text = header
            for row_idx, row_data in enumerate(rows):
                for col_idx, value in enumerate(row_data):
                    table.rows[row_idx + 1].cells[col_idx].text = str(value)
    
    doc.add_paragraph()
    footer = doc.add_paragraph(f"Erstellt mit AlSales AI • {datetime.now().strftime('%d.%m.%Y')}")
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()

